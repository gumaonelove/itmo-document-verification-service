"""Тесты парсинга документов (veridoc.parsing).

Сеть, LLM, эмбеддер и Qdrant не нужны — проверяется только parse_document и
его библиотеки (python-docx, python-pptx). Документы генерируются на лету во
временной директории (фикстура tmp_path).
"""

import pytest

from veridoc.parsing import parse_document


def test_parse_docx_paragraphs(tmp_path):
    """DOCX: пустой абзац отброшен, location=paragraph, тексты совпадают."""
    from docx import Document

    paragraphs = [
        "Выручка компании за 2023 год составила 150 млн рублей.",
        "Договор подписан 15 марта 2024 года директором Ивановым.",
        "",  # пустой абзац — должен быть отброшен
        "Штат компании насчитывает 42 сотрудника.",
    ]
    expected = [p for p in paragraphs if p]  # непустые, в порядке следования

    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    path = tmp_path / "отчёт.docx"
    doc.save(str(path))

    segments = parse_document(str(path))

    # Пустой абзац отброшен — остаются только непустые.
    assert len(segments) == len(expected)
    # doc_id — имя файла без расширения.
    assert all(seg.doc_id == "отчёт" for seg in segments)
    # location: kind=="paragraph" с целочисленным index.
    for seg in segments:
        assert seg.location["kind"] == "paragraph"
        assert isinstance(seg.location["index"], int)
    # Тексты совпадают и идут в исходном порядке.
    assert [seg.text for seg in segments] == expected


def test_parse_docx_index_skips_empty(tmp_path):
    """Сквозной индекс абзацев не сбрасывается на пустом абзаце."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("Первый непустой абзац.")
    doc.add_paragraph("")  # пустой между непустыми
    doc.add_paragraph("Третий непустой абзац.")
    path = tmp_path / "doc.docx"
    doc.save(str(path))

    segments = parse_document(str(path))

    assert len(segments) == 2
    # Индексы — 1-based и пропускают позицию пустого абзаца (1 и 3).
    indices = [seg.location["index"] for seg in segments]
    assert indices == [1, 3]
    # id детерминирован: "{doc_id}#paragraph{index}".
    assert segments[0].id == "doc#paragraph1"
    assert segments[1].id == "doc#paragraph3"


def test_parse_pptx_slides(tmp_path):
    """PPTX: два слайда, kind=="slide", index 1-based, тексты совпадают."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[5]  # макет "Только заголовок" — есть заголовок

    slide_texts = [
        "Квартальная прибыль выросла на 12 процентов.",
        "Запуск нового продукта запланирован на сентябрь.",
    ]
    for text in slide_texts:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = text

    path = tmp_path / "презентация.pptx"
    prs.save(str(path))

    segments = parse_document(str(path))

    # Два непустых слайда → два сегмента.
    assert len(segments) == 2
    assert all(seg.doc_id == "презентация" for seg in segments)
    # location: kind=="slide", index 1-based и по порядку.
    assert [seg.location["kind"] for seg in segments] == ["slide", "slide"]
    assert [seg.location["index"] for seg in segments] == [1, 2]
    # Текст каждого слайда совпадает с заданным.
    assert [seg.text for seg in segments] == slide_texts


def test_parse_unknown_extension_raises(tmp_path):
    """Неизвестное расширение → ValueError."""
    path = tmp_path / "файл.txt"
    path.write_text("просто текст", encoding="utf-8")

    with pytest.raises(ValueError):
        parse_document(str(path))
