"""Парсинг исходных документов в сегменты.

Единая точка входа — parse_document(path): по расширению выбирается парсер
(.pdf / .docx / .pptx) и документ режется на атомарные Segment с привязкой к
месту (страница / абзац / слайд). Пустые сегменты отбрасываются.

Идентификатор сегмента детерминирован: f"{doc_id}#{kind}{index}", где
doc_id — имя файла без расширения. Это даёт стабильные id между запусками.
"""

import logging
import os
import time

from veridoc.models import Segment

logger = logging.getLogger(__name__)


def parse_document(path: str) -> list[Segment]:
    """Распарсить документ в список сегментов.

    Диспетчеризация по расширению без учёта регистра: .pdf / .docx / .pptx.
    Иначе — ValueError. Возвращает сегменты в порядке следования в документе.
    """
    ext = os.path.splitext(path)[1].lower()
    doc_id = os.path.splitext(os.path.basename(path))[0]

    start = time.perf_counter()
    if ext == ".pdf":
        segments = _parse_pdf(path, doc_id)
    elif ext == ".docx":
        segments = _parse_docx(path, doc_id)
    elif ext == ".pptx":
        segments = _parse_pptx(path, doc_id)
    else:
        raise ValueError(f"Неподдерживаемое расширение: {ext!r} (ожидались .pdf/.docx/.pptx)")

    elapsed = time.perf_counter() - start
    logger.info(
        "Парсинг %s (%s): %d сегментов за %.2f c",
        os.path.basename(path), ext.lstrip("."), len(segments), elapsed,
    )
    return segments


def _segment(doc_id: str, kind: str, index: int, text: str) -> Segment:
    """Собрать Segment с детерминированным id и location."""
    return Segment(
        id=f"{doc_id}#{kind}{index}",
        doc_id=doc_id,
        text=text,
        location={"kind": kind, "index": index},
    )


def _parse_pdf(path: str, doc_id: str) -> list[Segment]:
    """PDF: по странице на сегмент.

    Основной путь — pymupdf4llm.to_markdown(page_chunks=True): для каждой
    страницы получаем словарь, текст берём из chunk["text"]. При недоступности
    или ошибке пакета — фолбэк на голый pymupdf (page.get_text()).
    """
    segments: list[Segment] = []
    try:
        import pymupdf4llm

        chunks = pymupdf4llm.to_markdown(path, page_chunks=True)
        for n, chunk in enumerate(chunks, start=1):
            text = (chunk.get("text") or "").strip()
            if text:
                segments.append(_segment(doc_id, "page", n, text))
        return segments
    except Exception as exc:  # фолбэк на pymupdf
        logger.warning("pymupdf4llm недоступен/упал (%s), фолбэк на pymupdf", exc)

    import pymupdf

    segments = []
    doc = pymupdf.open(path)
    try:
        for n, page in enumerate(doc, start=1):
            text = (page.get_text() or "").strip()
            if text:
                segments.append(_segment(doc_id, "page", n, text))
    finally:
        doc.close()
    return segments


def _parse_docx(path: str, doc_id: str) -> list[Segment]:
    """DOCX: по абзацу на сегмент, со сквозным счётчиком.

    Идём сначала по абзацам тела (doc.paragraphs), затем по ячейкам таблиц
    (doc.tables -> rows -> cells). paragraph-индекс единый и сквозной для
    абзацев и ячеек, чтобы id оставались уникальными.
    """
    from docx import Document

    doc = Document(path)
    segments: list[Segment] = []
    index = 0

    for paragraph in doc.paragraphs:
        index += 1
        text = (paragraph.text or "").strip()
        if text:
            segments.append(_segment(doc_id, "paragraph", index, text))

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                index += 1
                text = (cell.text or "").strip()
                if text:
                    segments.append(_segment(doc_id, "paragraph", index, text))

    return segments


def _parse_pptx(path: str, doc_id: str) -> list[Segment]:
    """PPTX: один сегмент на слайд.

    Текст слайда — конкатенация текста всех фигур, имеющих text_frame
    (1-based нумерация слайдов).
    """
    from pptx import Presentation

    prs = Presentation(path)
    segments: list[Segment] = []

    for n, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                frame_text = shape.text_frame.text.strip()
                if frame_text:
                    parts.append(frame_text)
        text = "\n".join(parts).strip()
        if text:
            segments.append(_segment(doc_id, "slide", n, text))

    return segments
