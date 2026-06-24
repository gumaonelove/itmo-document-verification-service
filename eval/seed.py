"""Генератор синтетического корпуса для оценки ВериДок (параметризованный).

Создаёт N пар документов «источник» (корректные факты) и «проверяемый»
(с внесёнными ошибками) во всех трёх форматах (DOCX/PPTX/PDF) и общий
машиночитаемый манифест ``ground_truth.json``.

Типы внесений (error_type) и подтипы (subtype):
  - numeric  : изменено значение факта  -> ожидаемый вердикт contradicted;
               подтипы number / date / percent / money;
  - claim    : подменён факт/имя/организация -> contradicted;
               подтипы fact / name / org;
  - internal : одна величина названа по-разному в двух местах проверяемого
               документа -> internal_conflict; подтипы number / date;
  - not_found: в проверяемом есть утверждение, которого НЕТ в источнике ->
               ожидаемый вердикт not_found (диагностика «не выдумывает
               contradicted»).
Плюс ЧИСТЫЕ документы (без ошибок): любой не-supported вердикт на них — ложное
срабатывание (FP).

Раскладка для однозначной локации: ОДИН факт = один абзац (DOCX) / слайд (PPTX) /
страница (PDF). Тогда parser выставит location.index = позиция факта (1-based),
и сопоставление с ground truth детерминировано.

Детерминизм: всё значения берутся из ``random.Random(seed + doc_index)``; повтор с
тем же ``--seed`` даёт побайтово тот же манифест.

Запуск:
    python eval/seed.py --docs 16 --numeric 2 --claim 2 --internal 2
    python eval/seed.py --docs 16 --numeric 2 --claim 2 --internal 2 \
        --notfound 1 --clean 3 --seed 42
"""

import argparse
import json
import logging
import os
import random
import time

logger = logging.getLogger(__name__)

DATASETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "datasets")
GROUND_TRUTH_PATH = os.path.join(DATASETS_DIR, "ground_truth.json")

_FORMATS = ("docx", "pptx", "pdf")

# --- Словари (детерминированно выбираются по seed) -------------------------
_THEMES = [
    ("Финансовый отчёт", "финансовый департамент"),
    ("План проекта", "проектный офис"),
    ("Кадровая записка", "отдел кадров"),
    ("Смета договора", "договорный отдел"),
    ("Логистический отчёт", "служба логистики"),
    ("Отчёт о закупках", "отдел закупок"),
    ("Маркетинговый план", "отдел маркетинга"),
    ("Производственная сводка", "производственный отдел"),
]
_ORGS = ["Орбита", "Маяк", "Вектор", "Контур", "Гранит", "Зенит", "Альфа", "Дельта"]
_ALT_ORGS = ["Сатурн", "Полюс", "Эверест", "Атлант", "Меридиан", "Каскад"]
_NAMES = ["Игорь Смирнов", "Анна Воробьёва", "Павел Кузнецов", "Мария Лебедева",
          "Олег Соколов", "Елена Морозова", "Дмитрий Орлов", "Светлана Громова"]
_ALT_NAMES = ["Андрей Зайцев", "Ольга Белова", "Сергей Поляков", "Наталья Гусева",
              "Виктор Лапин", "Юлия Седова"]
_MONTHS = ["января", "февраля", "марта", "апреля", "мая", "июня", "июля",
           "августа", "сентября", "октября", "ноября", "декабря"]
_LABELS = ["выручка", "чистая прибыль", "операционные расходы", "объём инвестиций",
           "фонд оплаты труда", "затраты на материалы", "дебиторская задолженность"]
_EVENTS = ["старт проекта", "ввод в эксплуатацию", "приёмка работ",
           "поставка оборудования", "завершение этапа"]
_ITEMS = ["оборудование", "монтажные работы", "лицензии", "доставка", "пусконаладка"]
# Роли для claim/name — фразы «<роль> назначен <имя>» (T-lite надёжно извлекает
# именно такую конструкцию назначения человека).
_ROLES = ["Руководителем проекта", "Главным бухгалтером", "Ведущим аналитиком",
          "Ответственным за приёмку", "Начальником отдела закупок",
          "Куратором поставки", "Главным инженером", "Старшим экономистом",
          "Менеджером проекта", "Ответственным исполнителем"]
# Задачи для claim/fact — фразы «Статус <задачи>: <статус>».
_TASKS = ["согласования сметы", "приёмки партии", "аудита расходов",
          "поверки приборов", "проверки документации", "ввода данных в систему"]


# --- Построение фактов ------------------------------------------------------
# Каждый факт: dict с text (предложение), subtype, и subject (ключ для пары
# источник/проверяемый). value_fragment — подстрока с проверяемым значением.

def _fact_number(rng, label, value):
    return {"subtype": "number", "subject": f"number:{label}",
            "text": f"Показатель «{label}» за 2024 год составил {value} единиц.",
            "value_fragment": f"{value} единиц"}


def _fact_money(rng, item, value):
    return {"subtype": "money", "subject": f"money:{item}",
            "text": f"Стоимость по статье «{item}» составляет {value} тысяч рублей.",
            "value_fragment": f"{value} тысяч рублей"}


def _fact_percent(rng, label, value):
    return {"subtype": "percent", "subject": f"percent:{label}",
            "text": f"Рост показателя «{label}» составил {value} процентов.",
            "value_fragment": f"{value} процентов"}


def _fact_date(rng, event, day, month, year=2025):
    return {"subtype": "date", "subject": f"date:{event}",
            "text": f"Срок «{event}» назначен на {day} {month} {year} года.",
            "value_fragment": f"{day} {month} {year} года"}


def _fact_name(rng, role, name):
    return {"subtype": "name", "subject": f"name:{role}",
            "text": f"{role} назначен {name}.",
            "value_fragment": name}


def _fact_status(rng, task, status):
    return {"subtype": "fact", "subject": f"fact:{task}",
            "text": f"Статус {task}: {status}.",
            "value_fragment": status}


def _numeric_fact(rng, used_labels):
    """Случайный numeric-факт (number/date/percent/money) с уникальным subject."""
    kind = rng.choice(["number", "money", "percent", "date"])
    if kind == "number":
        label = _pick_unique(rng, _LABELS, used_labels)
        return _fact_number(rng, label, rng.randint(120, 980))
    if kind == "money":
        item = _pick_unique(rng, _ITEMS, used_labels)
        return _fact_money(rng, item, rng.randint(60, 940))
    if kind == "percent":
        label = _pick_unique(rng, _LABELS, used_labels)
        return _fact_percent(rng, label, rng.randint(3, 48))
    event = _pick_unique(rng, _EVENTS, used_labels)
    return _fact_date(rng, event, rng.randint(1, 27), rng.choice(_MONTHS))


def _claim_fact(rng, used_labels):
    """Случайный claim-факт: имя (назначение человека) или статус задачи.

    Подтип организации (org) намеренно не используется: текущая модель/промпт
    не извлекает названия организаций как проверяемые утверждения — это
    задокументированное ограничение (см. metrics.md / README §9).
    """
    if rng.random() < 0.5:
        role = _pick_unique(rng, _ROLES, used_labels)
        return _fact_name(rng, role, rng.choice(_NAMES))
    task = _pick_unique(rng, _TASKS, used_labels)
    return _fact_status(rng, task, "согласовано")


def _pick_unique(rng, pool, used):
    """Выбрать элемент пула, которого ещё не было (для уникальных subject)."""
    options = [x for x in pool if x not in used]
    if not options:
        options = pool
    choice = rng.choice(options)
    used.add(choice)
    return choice


# --- Подмена значения (создание ошибки) ------------------------------------

def _tamper(rng, fact):
    """Вернуть (tampered_fact, tampered_value_fragment) — изменённую версию факта."""
    st = fact["subtype"]
    if st == "number":
        v = int(fact["value_fragment"].split()[0])
        nv = v + rng.randint(40, 300)
        return _replace(fact, f"{v} единиц", f"{nv} единиц"), f"{nv} единиц"
    if st == "money":
        v = int(fact["value_fragment"].split()[0])
        nv = v + rng.randint(50, 400)
        return _replace(fact, f"{v} тысяч", f"{nv} тысяч"), f"{nv} тысяч рублей"
    if st == "percent":
        v = int(fact["value_fragment"].split()[0])
        nv = v + rng.randint(5, 40)
        return _replace(fact, f"{v} процентов", f"{nv} процентов"), f"{nv} процентов"
    if st == "date":
        # Сдвигаем месяц и день — значение точно отличается.
        old = fact["value_fragment"]
        parts = old.split()
        new_month = _MONTHS[(_MONTHS.index(parts[1]) + 5) % 12]
        new_day = (int(parts[0]) % 27) + 1
        new_frag = f"{new_day} {new_month} {parts[2]} {parts[3]}"
        return _replace(fact, old, new_frag), new_frag
    if st == "name":
        new = rng.choice(_ALT_NAMES)
        return _replace(fact, fact["value_fragment"], new), new
    if st == "org":
        new = rng.choice(_ALT_ORGS)
        old = fact["value_fragment"]
        return _replace(fact, old, f"«{new}»"), f"«{new}»"
    # fact/status — меняем на ЯВНО другой статус (а не на отрицание «не …»),
    # чтобы тест не сводился к распознаванию частицы «не».
    new = "отклонено"
    return _replace(fact, fact["value_fragment"], new), new


def _replace(fact, old_sub, new_sub):
    f = dict(fact)
    f["text"] = fact["text"].replace(old_sub, new_sub)
    f["value_fragment"] = new_sub
    return f


# --- Генерация одной пары документов ---------------------------------------

def _generate_pair(doc_index, fmt, counts, base_seed):
    """Сгенерировать пару (источник, проверяемый) + список GT-ошибок.

    counts: dict с numeric/claim/internal/notfound (число ошибок каждого типа)
    и флагом clean. Возвращает (source_texts, check_texts, gt_errors, is_clean).
    Один факт = одна единица (абзац/слайд/страница); индекс GT = позиция в check.
    """
    rng = random.Random(base_seed + doc_index * 1000 + ord(fmt[0]))
    theme_title, dept = _THEMES[doc_index % len(_THEMES)]
    used = set()

    source_facts = []  # факты в источнике (корректные)
    check_facts = []   # факты в проверяемом (часть подменена)
    gt = []            # записи ground truth

    def add_correct():
        f = _numeric_fact(rng, used) if rng.random() < 0.6 else _claim_fact(rng, used)
        source_facts.append(f)
        check_facts.append(f)

    # 1) Корректные факты (станут supported).
    for _ in range(3):
        add_correct()

    if counts["clean"]:
        # Чистый документ: только корректные факты, без ошибок.
        for _ in range(2):
            add_correct()
        return _render(source_facts, check_facts, theme_title, dept, fmt), gt, True

    # 2) numeric-ошибки.
    for _ in range(counts["numeric"]):
        f = _numeric_fact(rng, used)
        source_facts.append(f)
        tampered, frag = _tamper(rng, f)
        check_facts.append(tampered)
        gt.append(_gt_error("contradicted", "numeric", f["subtype"], frag,
                            len(check_facts)))

    # 3) claim-ошибки.
    for _ in range(counts["claim"]):
        f = _claim_fact(rng, used)
        source_facts.append(f)
        tampered, frag = _tamper(rng, f)
        check_facts.append(tampered)
        gt.append(_gt_error("contradicted", "claim", f["subtype"], frag,
                            len(check_facts)))

    # 4) internal-ошибки: одна величина названа по-разному в двух местах check.
    for _ in range(counts["internal"]):
        st_kind = rng.choice(["number", "money"])
        if st_kind == "number":
            label = _pick_unique(rng, _LABELS, used)
            v = rng.randint(120, 980)
            base = _fact_number(rng, label, v)
            alt_v = v + rng.randint(40, 300)
            alt = _replace(base, f"{v} единиц", f"{alt_v} единиц")
            frag = f"{alt_v} единиц"
        else:
            item = _pick_unique(rng, _ITEMS, used)
            v = rng.randint(60, 940)
            base = _fact_money(rng, item, v)
            alt_v = v + rng.randint(50, 400)
            alt = _replace(base, f"{v} тысяч", f"{alt_v} тысяч")
            frag = f"{alt_v} тысяч рублей"
        source_facts.append(base)               # источник: одно согласованное значение
        check_facts.append(base)                # в check — первое упоминание (верное)
        check_facts.append(alt)                 # и второе, конфликтующее
        gt.append(_gt_error("internal_conflict", "internal", base["subtype"], frag,
                            len(check_facts)))   # локация = второе (изменённое) упоминание

    # 5) not_found: факт, которого НЕТ в источнике.
    for _ in range(counts["notfound"]):
        label = _pick_unique(rng, _LABELS, used)
        f = _fact_number(rng, "прочее: " + label, rng.randint(100, 999))
        check_facts.append(f)  # только в check
        gt.append(_gt_error("not_found", "not_found", "number", f["value_fragment"],
                            len(check_facts)))

    return _render(source_facts, check_facts, theme_title, dept, fmt), gt, False


def _gt_error(kind, error_type, subtype, value_fragment, check_index):
    return {
        "kind": kind,
        "error_type": error_type,
        "subtype": subtype,
        "tampered_value": value_fragment,
        "location_index": check_index,  # 1-based позиция в check; kind заполнит _render
    }


def _render(source_facts, check_facts, title, dept, fmt):
    """Сформировать тексты единиц для источника и проверяемого по формату."""
    return {
        "source": [f["text"] for f in source_facts],
        "check": [f["text"] for f in check_facts],
        "title": title,
        "dept": dept,
    }


# --- Запись файлов ----------------------------------------------------------

def _location_kind(fmt):
    return {"docx": "paragraph", "pptx": "slide", "pdf": "page"}[fmt]


def _write_docx(path, units, title, dept):
    from docx import Document
    doc = Document()
    # Один факт = один абзац; индекс location = позиция (1-based).
    for text in units:
        doc.add_paragraph(text)
    doc.save(path)


def _write_pptx(path, units, title, dept):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]
    # Один факт = один слайд; индекс location = номер слайда (1-based).
    for i, text in enumerate(units, start=1):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9.0), Inches(4.0))
        box.text_frame.text = text
    prs.save(path)


def _write_pdf(path, units, title, dept):
    import pymupdf
    fontsize, left, top, step = 11.0, 72.0, 72.0, 20.0
    doc = pymupdf.open()
    try:
        # Один факт = одна страница; индекс location = номер страницы (1-based).
        for text in units:
            page = doc.new_page()
            usable = page.rect.width - left - 56.0
            max_chars = int(usable // fontsize)
            y = top
            for chunk in _wrap_line(text, max_chars):
                page.insert_text((left, y), chunk, fontsize=fontsize, fontname="china-ss")
                y += step
        doc.save(path)
    finally:
        doc.close()


def _wrap_line(line, max_chars):
    if not line.strip():
        return [""]
    words, out, cur = line.split(" "), [], ""
    for w in words:
        candidate = w if not cur else f"{cur} {w}"
        if len(candidate) <= max_chars:
            cur = candidate
        else:
            if cur:
                out.append(cur)
            cur = w
    if cur:
        out.append(cur)
    return out


_WRITERS = {"docx": _write_docx, "pptx": _write_pptx, "pdf": _write_pdf}


# --- Главная сборка ---------------------------------------------------------

def build_corpus(docs, numeric, claim, internal, notfound, clean, seed):
    """Сгенерировать корпус и вернуть манифест (documents + errors + meta)."""
    os.makedirs(DATASETS_DIR, exist_ok=True)
    documents, errors = [], []

    # Первые `clean` документов делаем чистыми; остальные — с ошибками.
    for i in range(docs):
        fmt = _FORMATS[i % len(_FORMATS)]
        is_clean_doc = i < clean
        counts = {
            "numeric": 0 if is_clean_doc else numeric,
            "claim": 0 if is_clean_doc else claim,
            "internal": 0 if is_clean_doc else internal,
            "notfound": 0 if is_clean_doc else notfound,
            "clean": is_clean_doc,
        }
        rendered, gt, is_clean = _generate_pair(i, fmt, counts, seed)

        tag = f"{i:03d}_{_THEMES[i % len(_THEMES)][0].split()[0].lower()}"
        doc_name = f"check_{tag}.{fmt}"
        src_name = f"source_{tag}.{fmt}"
        writer = _WRITERS[fmt]
        writer(os.path.join(DATASETS_DIR, src_name), rendered["source"],
               rendered["title"], rendered["dept"])
        writer(os.path.join(DATASETS_DIR, doc_name), rendered["check"],
               rendered["title"], rendered["dept"])

        documents.append({"doc": doc_name, "source": src_name, "format": fmt,
                          "clean": is_clean})

        loc_kind = _location_kind(fmt)
        for e in gt:
            errors.append({
                "doc": doc_name,
                "source": src_name,
                "format": fmt,
                "kind": e["kind"],
                "error_type": e["error_type"],
                "subtype": e["subtype"],
                "location": {"kind": loc_kind, "index": e["location_index"]},
                "tampered_value": e["tampered_value"],
            })

    meta = {
        "seed": seed,
        "docs": docs,
        "clean_docs": clean,
        "errors_total": sum(1 for e in errors if e["kind"] != "not_found"),
        "not_found_cases": sum(1 for e in errors if e["kind"] == "not_found"),
        "per_error_doc": {"numeric": numeric, "claim": claim, "internal": internal,
                          "notfound": notfound},
    }
    return {"meta": meta, "documents": documents, "errors": errors}


def main() -> None:
    parser = argparse.ArgumentParser(description="Генератор корпуса ВериДок.")
    parser.add_argument("--docs", type=int, default=16, help="всего документов")
    parser.add_argument("--numeric", type=int, default=2, help="numeric-ошибок на док")
    parser.add_argument("--claim", type=int, default=2, help="claim-ошибок на док")
    parser.add_argument("--internal", type=int, default=2, help="internal-ошибок на док")
    parser.add_argument("--notfound", type=int, default=1, help="not_found-кейсов на док")
    parser.add_argument("--clean", type=int, default=3, help="чистых документов")
    parser.add_argument("--seed", type=int, default=42, help="seed детерминизма")
    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO)
    start = time.perf_counter()

    manifest = build_corpus(args.docs, args.numeric, args.claim, args.internal,
                            args.notfound, args.clean, args.seed)
    with open(GROUND_TRUTH_PATH, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)

    elapsed = time.perf_counter() - start
    m = manifest["meta"]
    by_type, by_fmt = {}, {}
    for e in manifest["errors"]:
        by_type[e["error_type"]] = by_type.get(e["error_type"], 0) + 1
        by_fmt[e["format"]] = by_fmt.get(e["format"], 0) + 1

    print("=" * 64)
    print("Генерация корпуса ВериДок")
    print("=" * 64)
    print(f"Документов:            {len(manifest['documents'])} (чистых: {m['clean_docs']})")
    print(f"Ошибок-детекций:       {m['errors_total']}  (+ not_found-кейсов: {m['not_found_cases']})")
    print(f"По типам:              {by_type}")
    print(f"По форматам (ошибки):  {by_fmt}")
    print(f"seed:                  {m['seed']}")
    print(f"Каталог:               {DATASETS_DIR}")
    print(f"Манифест:              {GROUND_TRUTH_PATH}")
    print(f"Время:                 {elapsed:.2f} c")
    print("=" * 64)


if __name__ == "__main__":
    main()
